from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_25_ui_ux.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Frontend және UI/UX: Қолданушы ыңғайлылығы"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. СОЛ ЖАҚ: МӘТІН (Логика)
# ============================================================
left_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.0)
)
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = RGBColor(248, 249, 249)
left_bg.line.color.rgb = RGBColor(52, 152, 219)
left_bg.line.width = Pt(2)

txLeft = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.6), Inches(4.8))
tL = txLeft.text_frame
tL.word_wrap = True

def add_bullet(tf, icon, title, desc):
    pT = tf.add_paragraph()
    pT.text = f"{icon} {title}"
    pT.font.bold = True
    pT.font.size = Pt(18)
    pT.font.color.rgb = RGBColor(41, 128, 185)
    pT.space_before = Pt(10)
    
    pD = tf.add_paragraph()
    pD.text = desc
    pD.font.size = Pt(14)
    pD.font.color.rgb = RGBColor(52, 73, 94)

p0 = tL.paragraphs[0]
p0.text = "Неге React және Dark Mode?"
p0.font.bold = True
p0.font.size = Pt(22)
p0.font.color.rgb = RGBColor(44, 62, 80)

add_bullet(tL, "⚛️", "Жылдамдық (React Vite)", "Virtual DOM арқылы интерфейс қатып қалмайды. Таймер мен AI дәлдігі (Accuracy) миллисекунд ішінде жаңартылып отырады.")
add_bullet(tL, "👁️", "Қолданушы көзі (Контраст)", "Адам камерадан 2 метр алыс тұрғанда да көруі үшін Dark Mode (қараңғы фон) және өте үлкен ашық түсті мәтіндер қолданылды.")
add_bullet(tL, "👆", "Қолданушы саусағы (Mobile-first)", "Жаттығу кезінде терлеген қолмен басуға ыңғайлы болуы үшін UI-да майда детальдар емес, үп-үлкен Карточкалар (Cards) мен Түймелер жасалды.")

# ============================================================
# 3. ОҢ ЖАҚ: СУРЕТТЕР ОРНЫ (Mockups)
# ============================================================
# Since we cannot automatically load chat images via python script easily, 
# we create beautiful placeholders with instructions. 
# The user can just drag and drop the images over these shapes.

def create_image_placeholder(top, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.6), Inches(top), Inches(7.2), Inches(2.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(34, 40, 49) # Dark UI color to match their screenshot
    shape.line.color.rgb = RGBColor(241, 196, 15)
    shape.line.width = Pt(2)
    shape.line.dash_style = 4 # Dashed to indicate placeholder
    
    tx = slide.shapes.add_textbox(Inches(5.6), Inches(top + 0.9), Inches(7.2), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

create_image_placeholder(1.2, "ОСЫ ЖЕРГЕ 'БАСТЫ БЕТ' СКРИНШОТЫН ҚОЙЫҢЫЗ\n(Drag and drop here)")
create_image_placeholder(3.8, "ОСЫ ЖЕРГЕ 'ЖАТТЫҒУЛАР' СКРИНШОТЫН ҚОЙЫҢЫЗ\n(Drag and drop here)")

# ============================================================
# 4. ТӨМЕНДЕГІ ТҮСІНІКТЕМЕ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.33), Inches(1.0))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Қорытынды: Код қанша мықты болса да, интерфейс ыңғайсыз болса — жоба сәтсіз. "
    "Сондықтан біз UI/UX дизайнын спортпен шұғылданатын адамның нақты анатомиясы "
    "мен қимылына (алыстан көру, үлкен түймелер басу) бейімдедік."
)
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(16)
pFoot.font.color.rgb = RGBColor(142, 68, 173)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
