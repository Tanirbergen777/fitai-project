from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_data_preprocessing.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Data Preprocessing: Мақсатты Аудитория және Деректерді Тазалау"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ФИЛЬТР СХЕМАСЫ (Сол жақта)
# ============================================================
# Бастапқы база
start_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.5), Inches(1.0))
start_box.fill.solid()
start_box.fill.fore_color.rgb = RGBColor(149, 165, 166)
tx1 = start_box.text_frame
tx1.paragraphs[0].text = "Шикі Датасет (NHANES)\nБалалар + Ересектер + Қарттар"
tx1.paragraphs[0].font.size = Pt(16)
tx1.paragraphs[0].alignment = PP_ALIGN.CENTER

# Сүзу (Фильтр)
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.3), Inches(2.9), Inches(0.5), Inches(0.5)).fill.solid()

filter_box = slide.shapes.add_shape(MSO_SHAPE.SNIP_2_DIAG_RECTANGLE, Inches(0.8), Inches(3.5), Inches(3.5), Inches(1.5))
filter_box.fill.solid()
filter_box.fill.fore_color.rgb = RGBColor(192, 57, 43)
tx2 = filter_box.text_frame
tx2.paragraphs[0].text = "ДЕРЕКТЕРДІ ТАЗАЛАУ (Алып тасталды):"
tx2.paragraphs[0].font.bold = True
tx2.paragraphs[0].font.size = Pt(14)
tx2.paragraphs[0].alignment = PP_ALIGN.CENTER
p2 = tx2.add_paragraph()
p2.text = "❌ 18 жасқа дейінгі балалар\n❌ Экстремалды салмақ (Outliers)\n❌ Толық емес мәліметтер"
p2.font.size = Pt(14)

# Нәтиже
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.3), Inches(5.1), Inches(0.5), Inches(0.5)).fill.solid()

end_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.7), Inches(3.5), Inches(1.0))
end_box.fill.solid()
end_box.fill.fore_color.rgb = RGBColor(39, 174, 96)
tx3 = end_box.text_frame
tx3.paragraphs[0].text = "Таза Датасет (Мақсатты Аудитория)\nТек 18 - 80 жас аралығы (Ересектер)"
tx3.paragraphs[0].font.bold = True
tx3.paragraphs[0].font.size = Pt(16)
tx3.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# 3. ТҮСІНІКТЕМЕ ЖӘНЕ ГРАФИК ИМИТАЦИЯСЫ (Оң жақта)
# ============================================================
# График фоны
graph_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.8), Inches(7.5), Inches(2.8))
graph_bg.fill.solid()
graph_bg.fill.fore_color.rgb = RGBColor(240, 243, 244)
graph_bg.line.color.rgb = RGBColor(189, 195, 199)

# Имитация гистограммы (Жас ерекшеліктері)
ages = ["18-25", "26-35", "36-45", "46-55", "56-65", "66-80"]
heights = [1.5, 2.2, 2.0, 1.8, 1.3, 0.8]
colors = [(52, 152, 219)] * 6

for i, h in enumerate(heights):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5 + i*1.1), Inches(4.5 - h), Inches(0.6), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*colors[i])
    bar.line.color.rgb = RGBColor(41, 128, 185)
    
    # Label under bar
    txL = slide.shapes.add_textbox(Inches(5.4 + i*1.1), Inches(4.5), Inches(0.8), Inches(0.3))
    pL = txL.text_frame.paragraphs[0]
    pL.text = ages[i]
    pL.font.size = Pt(12)
    pL.alignment = PP_ALIGN.CENTER

# Түсініктеме мәтіні
txRight = slide.shapes.add_textbox(Inches(5.0), Inches(5.0), Inches(7.5), Inches(1.8))
txR = txRight.text_frame
txR.word_wrap = True

pR1 = txR.paragraphs[0]
pR1.text = "Неліктен балалар алынып тасталды?"
pR1.font.bold = True
pR1.font.size = Pt(20)
pR1.font.color.rgb = RGBColor(44, 62, 80)

pR2 = txR.add_paragraph()
pR2.text = "Балалардың метаболизмі мен дене құрылысы (өсу процесі) ересектерден мүлде бөлек жұмыс істейді. Егер оларды датасетке қоссақ, Machine Learning моделі шатасып, ересек адамға баланың жаттығуын беріп қоюы мүмкін."
pR2.font.size = Pt(16)
pR2.font.color.rgb = RGBColor(52, 73, 94)
pR2.space_before = Pt(8)

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
