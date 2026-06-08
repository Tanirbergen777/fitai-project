from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_27_swot_analysis.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Стратегиялық бағалау: SWOT Анализ"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. SWOT МАТРИЦАСЫ (2x2)
# ============================================================
def create_swot_card(x, y, w, h, icon, title, bullets, main_color, bg_color):
    # Background card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*bg_color)
    card.line.color.rgb = RGBColor(*main_color)
    card.line.width = Pt(2)
    
    # Title Header Block
    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*main_color)
    header.line.fill.background()
    
    txH = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(w), Inches(0.6))
    pH = txH.text_frame.paragraphs[0]
    pH.text = f"{icon}  {title}"
    pH.font.bold = True
    pH.font.size = Pt(22)
    pH.font.color.rgb = RGBColor(255, 255, 255)
    pH.alignment = PP_ALIGN.CENTER
    
    # Bullets Content
    txC = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9), Inches(w - 0.4), Inches(h - 1.0))
    tC = txC.text_frame
    tC.word_wrap = True
    
    for idx, bullet in enumerate(bullets):
        pC = tC.add_paragraph() if idx > 0 else tC.paragraphs[0]
        pC.text = f"• {bullet}"
        pC.font.size = Pt(15)
        pC.font.color.rgb = RGBColor(44, 62, 80)
        pC.space_before = Pt(8)

w, h = 5.8, 2.8
x_left = 0.6
x_right = 6.9
y_top = 1.3
y_bottom = 4.3

# S: Strengths (Top-Left)
s_bullets = [
    "AI арқылы 100% жекелендірілген тәжірибе (Personalization).",
    "Computer Vision арқылы нақты уақытта (Real-time) қатені түзету.",
    "Масштабтауға (scale) дайын таза микросервистік архитектура."
]
create_swot_card(x_left, y_top, w, h, "💪", "Күшті жақтар (S)", s_bullets, (39, 174, 96), (232, 248, 245))

# W: Weaknesses (Top-Right)
w_bullets = [
    "Әзірге CV модельдері тек 3 жаттығуды (Squat, Push-up, Lunge) ғана таниды.",
    "Камераның дұрыс жұмыс істеуі үшін жақсы жарықтандыруды қажет етеді."
]
create_swot_card(x_right, y_top, w, h, "📉", "Кемшіліктер (W)", w_bullets, (192, 57, 43), (253, 237, 236))

# O: Opportunities (Bottom-Left)
o_bullets = [
    "Мобильді қосымша (React Native) ретінде шығарып, жаһандық нарыққа шығу.",
    "B2B нарығы: Фитнес клубтарға немесе оңалту (Rehab) орталықтарына API сату."
]
create_swot_card(x_left, y_bottom, w, h, "🚀", "Мүмкіндіктер (O)", o_bullets, (41, 128, 185), (235, 245, 251))

# T: Threats (Bottom-Right)
t_bullets = [
    "Нарықтағы үлкен фитнес қолданбалардың (Freeletics, Fitbod) бәсекелестігі.",
    "Қолданушылар көбейген сайын серверлік (LLM, GPU) шығындардың күрт артуы."
]
create_swot_card(x_right, y_bottom, w, h, "⚠️", "Қауіп-қатерлер (T)", t_bullets, (211, 84, 0), (253, 235, 208))

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
