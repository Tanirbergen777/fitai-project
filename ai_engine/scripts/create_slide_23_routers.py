from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_23_routers.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Модульдік Backend: 7 Router және 30+ Endpoint"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. СОЛ ЖАҚ: МӘТІНДІК СИПАТТАМА
# ============================================================
txLeft = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
tL = txLeft.text_frame
tL.word_wrap = True

def add_bullet(tf, title, desc):
    pT = tf.add_paragraph()
    pT.text = title
    pT.font.bold = True
    pT.font.size = Pt(20)
    pT.font.color.rgb = RGBColor(41, 128, 185)
    pT.space_before = Pt(14)
    
    pD = tf.add_paragraph()
    pD.text = desc
    pD.font.size = Pt(16)
    pD.font.color.rgb = RGBColor(52, 73, 94)

p0 = tL.paragraphs[0]
p0.text = "Модульдік Архитектура:"
p0.font.bold = True
p0.font.size = Pt(22)
p0.font.color.rgb = RGBColor(44, 62, 80)

add_bullet(tL, "• Монолит емес (No Monolith)", "Біз барлық кодты бір файлға (main.py) тықпадық. Оның орнына әр қызметті жеке Router-ге бөлдік. Бұл жобаны масштабтауды (scale) өте жеңілдетеді.")
add_bullet(tL, "• Мақсаттардың бөлінуі", "AI логикасы жеке, деректер базасы жеке, пайдаланушы профилі жеке папкаларда (Separation of Concerns) орналасқан.")
add_bullet(tL, "• Swagger UI", "FastAPI барлық 30+ endpoint үшін интерактивті құжаттаманы (Swagger) автоматты түрде жасап берді.")

# ============================================================
# 3. ОҢ ЖАҚ: СХЕМА (ROUTERS)
# ============================================================
# 3.1 Орталық блок (FastAPI Hub)
hub_x, hub_y, hub_w, hub_h = 7.5, 1.2, 3.5, 1.0
hub_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(hub_x), Inches(hub_y), Inches(hub_w), Inches(hub_h)
)
hub_shape.fill.solid()
hub_shape.fill.fore_color.rgb = RGBColor(39, 174, 96)
hub_shape.line.color.rgb = RGBColor(25, 111, 61)
hub_shape.line.width = Pt(2)

txHub = slide.shapes.add_textbox(Inches(hub_x), Inches(hub_y + 0.1), Inches(hub_w), Inches(hub_h))
tH = txHub.text_frame
tH.word_wrap = True
pH1 = tH.paragraphs[0]
pH1.text = "⚙️ FastAPI Application"
pH1.font.bold = True
pH1.font.size = Pt(20)
pH1.font.color.rgb = RGBColor(255, 255, 255)
pH1.alignment = PP_ALIGN.CENTER
pH2 = tH.add_paragraph()
pH2.text = "Main Entry Point (api/main.py)"
pH2.font.size = Pt(14)
pH2.font.color.rgb = RGBColor(234, 250, 241)
pH2.alignment = PP_ALIGN.CENTER

# 3.2 7 Router Блоктары
routers = [
    {"icon": "👤", "path": "/users", "desc": "Тіркелу, профиль (Auth)"},
    {"icon": "🍎", "path": "/nutrition", "desc": "Тамақтану (AI) сақтау"},
    {"icon": "🏋️", "path": "/workouts", "desc": "Жаттығу тарихы"},
    {"icon": "📷", "path": "/vision", "desc": "3D нүктелер (MediaPipe)"},
    {"icon": "🤖", "path": "/chat", "desc": "Groq Llama 3.3 кеңесші"},
    {"icon": "📈", "path": "/analytics", "desc": "Прогресс, салмақ"},
    {"icon": "🔥", "path": "/motivation", "desc": "Рейтинг, от (streak)"}
]

# Grid parameters
start_x = 5.2
start_y = 3.2
card_w = 3.8
card_h = 0.9
x_spacing = 4.2
y_spacing = 1.1

def create_router_card(cx, cy, cw, ch, r_data):
    r_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx), Inches(cy), Inches(cw), Inches(ch)
    )
    r_shape.fill.solid()
    r_shape.fill.fore_color.rgb = RGBColor(235, 245, 251)
    r_shape.line.color.rgb = RGBColor(41, 128, 185)
    r_shape.line.width = Pt(1.5)

    tx = slide.shapes.add_textbox(Inches(cx), Inches(cy + 0.05), Inches(cw), Inches(ch))
    t = tx.text_frame
    t.word_wrap = True
    
    p1 = t.paragraphs[0]
    p1.text = f"{r_data['icon']} {r_data['path']}"
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(41, 128, 185)
    
    p2 = t.add_paragraph()
    p2.text = r_data['desc']
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(52, 73, 94)
    
    # Draw line from Hub to this card
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        Inches(hub_x + hub_w/2), Inches(hub_y + hub_h),
        Inches(cx + cw/2), Inches(cy)
    )
    connector.line.width = Pt(1.5)
    connector.line.color.rgb = RGBColor(149, 165, 166)

# Create 7 cards in a visually appealing arrangement
# Left column (4 cards)
for i in range(4):
    create_router_card(start_x, start_y + i * y_spacing, card_w, card_h, routers[i])

# Right column (3 cards)
for i in range(4, 7):
    # offset right column slightly down so it looks balanced (4 vs 3)
    create_router_card(start_x + x_spacing, start_y + 0.5 + (i - 4) * y_spacing, card_w, card_h, routers[i])

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
