from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_21_tech_stack.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жобаның Технологиялық Стекі және Жалпы Схемасы"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ДИАГРАММА / СХЕМА
# ============================================================

def create_card(left, top, width, height, title, subtitle, bg_color, border_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = Pt(3)

    tx = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(width), Inches(height))
    t = tx.text_frame
    t.word_wrap = True
    
    p1 = t.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(*border_color)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = t.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(52, 73, 94)
    p2.alignment = PP_ALIGN.CENTER
    
    return shape

# --- FRONTEND ---
front = create_card(
    1.0, 3.0, 2.5, 1.8, 
    "💻 Frontend", "React (Vite)\nTypeScript\nUI/UX", 
    (235, 245, 251), (41, 128, 185)
)

# --- BACKEND ---
# Backend is hosted on Render, so let's make a big container for Backend + Render
render_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.5), Inches(2.2), Inches(4.0), Inches(3.4)
)
render_bg.fill.solid()
render_bg.fill.fore_color.rgb = RGBColor(244, 246, 247)
render_bg.line.color.rgb = RGBColor(149, 165, 166)
render_bg.line.width = Pt(1.5)
render_bg.line.dash_style = 4  # dashed

# Render text
txR = slide.shapes.add_textbox(Inches(4.5), Inches(5.2), Inches(4.0), Inches(0.4))
tR = txR.text_frame
pR = tR.paragraphs[0]
pR.text = "☁️ Hosted on Render"
pR.font.size = Pt(14)
pR.font.italic = True
pR.font.color.rgb = RGBColor(127, 140, 141)
pR.alignment = PP_ALIGN.CENTER

back = create_card(
    5.0, 2.6, 3.0, 2.4, 
    "⚙️ Backend", "Python + FastAPI\nMachine Learning\nGroq LLM API", 
    (232, 248, 245), (39, 174, 96)
)

# --- DATABASE (NEON) ---
neon = create_card(
    9.5, 1.8, 2.8, 1.6, 
    "🗄️ Database", "Neon (PostgreSQL)\nRelational Data", 
    (253, 235, 208), (211, 84, 0)
)

# --- AUTH & STORAGE (SUPABASE) ---
supa = create_card(
    9.5, 4.2, 2.8, 1.6, 
    "🔒 Auth & Cloud", "Supabase\nAuthentication\nFile Storage", 
    (232, 218, 239), (142, 68, 173)
)

# ============================================================
# 3. СТРЕЛКАЛАР ЖӘНЕ БАЙЛАНЫС
# ============================================================
def add_arrow(x1, y1, x2, y2, text):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.width = Pt(3)
    connector.line.color.rgb = RGBColor(149, 165, 166)
    
    # Arrow head
    # Cannot easily set end_arrow_head via simple API, but line itself is okay, 
    # we can use a small shape or just text to show direction.
    
    if text:
        # text box near the line
        tx = slide.shapes.add_textbox(Inches((x1+x2)/2 - 0.5), Inches((y1+y2)/2 - 0.2), Inches(1.5), Inches(0.4))
        t = tx.text_frame
        p = t.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.alignment = PP_ALIGN.CENTER

# Front -> Back
add_arrow(3.5, 3.9, 5.0, 3.9, "REST API ➔")

# Back -> Neon
add_arrow(8.0, 2.6, 9.5, 2.6, "SQL ➔")

# Back -> Supabase
add_arrow(8.0, 5.0, 9.5, 5.0, "API ➔")

# Front -> Supabase (direct auth)
add_arrow(2.25, 4.8, 9.5, 5.5, "Auth / Media ➔")


# ============================================================
# 4. ТӨМЕНДЕГІ ТҮСІНІКТЕМЕ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.33), Inches(1.0))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Архитектуралық ерекшелік: Жоба локальді компьютерде емес, толығымен "
    "бұлттық (Cloud) технологиялар экожүйесінде (Render, Neon, Supabase) "
    "нақты өндіріске (Production) дайындалған."
)
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(16)
pFoot.font.color.rgb = RGBColor(52, 73, 94)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
