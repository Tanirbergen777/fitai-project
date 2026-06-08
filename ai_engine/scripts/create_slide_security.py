from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_security_stability.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жүйенің Қауіпсіздігі (Security) мен Тұрақтылығы"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ЕКІ БАҒАНДЫ КАРТОЧКАЛАР
# ============================================================
def create_info_card(x, y, w, h, icon, title, bullets, main_color, bg_color):
    # Background card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*bg_color)
    card.line.color.rgb = RGBColor(*main_color)
    card.line.width = Pt(2)
    
    # Header
    txH = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(w), Inches(0.6))
    pH = txH.text_frame.paragraphs[0]
    pH.text = f"{icon}  {title}"
    pH.font.bold = True
    pH.font.size = Pt(24)
    pH.font.color.rgb = RGBColor(*main_color)
    pH.alignment = PP_ALIGN.CENTER
    
    # Bullets Content
    txC = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.0), Inches(w - 0.6), Inches(h - 1.2))
    tC = txC.text_frame
    tC.word_wrap = True
    
    for idx, (b_title, b_desc) in enumerate(bullets):
        pT = tC.add_paragraph() if idx > 0 else tC.paragraphs[0]
        pT.text = f"• {b_title}"
        pT.font.bold = True
        pT.font.size = Pt(18)
        pT.font.color.rgb = RGBColor(44, 62, 80)
        pT.space_before = Pt(14) if idx > 0 else Pt(0)
        
        pD = tC.add_paragraph()
        pD.text = f"  {b_desc}"
        pD.font.size = Pt(15)
        pD.font.color.rgb = RGBColor(52, 73, 94)

w, h = 5.8, 4.5
x_left = 0.6
x_right = 6.9
y_pos = 1.6

# SECURITY CARD (Left)
security_bullets = [
    ("Supabase & JWT Авторизациясы:", "Құпия сөздер ашық сақталмайды. Пайдаланушылар қауіпсіз шифрланған JWT токендер арқылы қорғалған."),
    ("Pydantic Validation (FastAPI):", "Сырттан келген барлық деректер автоматты түрде тексеріледі. Бұл жүйені SQL Injection және XSS хакерлік шабуылдарынан қорғайды.")
]
create_info_card(x_left, y_pos, w, h, "🔒", "Деректер Қауіпсіздігі (Security)", security_bullets, (192, 57, 43), (253, 237, 236))

# STABILITY CARD (Right)
stability_bullets = [
    ("Асинхронды архитектура (Async):", "FastAPI-дің арқасында 1000 адам бірден кірсе де, сервер қатып қалмай сұраныстарды параллельді өңдейді."),
    ("Модульдік оқшаулау (Isolation):", "Код микросервистік рутерлерге бөлінген. Бір жерде баг (қате) шықса, ол бүкіл қосымшаны құлатпайды."),
    ("Serverless Database (Neon):", "Жүктеме артқан кезде деректер базасы автоматты түрде өз күшін ұлғайтады (Auto-scaling).")
]
create_info_card(x_right, y_pos, w, h, "⚙️", "Жүйе Тұрақтылығы (Scalability)", stability_bullets, (39, 174, 96), (232, 248, 245))

# ============================================================
# 3. ҚОРЫТЫНДЫ (Highlight)
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "Архитектура басынан бастап масштабтауға және жоғары жүктемелерге (High Load) төтеп беруге арналып жасалған."
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(18)
pFoot.font.color.rgb = RGBColor(41, 128, 185)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
