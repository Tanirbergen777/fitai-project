from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_22_backend_choice.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Бланк слайд қосу
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Backend: Неге дәл Python және FastAPI таңдалды?"
p.font.bold = True
p.font.size = Pt(34)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. КРЕАТИВТІ КАРТОЧКАЛАР
# ============================================================

def create_vs_card(left, top, icon, main_tech, not_tech, reason_title, reason_text, main_color, accent_color):
    # Негізгі фон (Card Background)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(5.8), Inches(4.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 249, 249)
    bg.line.color.rgb = RGBColor(*main_color)
    bg.line.width = Pt(3)

    # Жоғарғы Header (Түсті блок)
    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, # We'll overlap it slightly to look like a header
        Inches(left), Inches(top), Inches(5.8), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*main_color)
    header.line.fill.background()

    txH = slide.shapes.add_textbox(Inches(left), Inches(top + 0.2), Inches(5.8), Inches(1.0))
    tH = txH.text_frame
    pH = tH.paragraphs[0]
    pH.text = f"{icon} {main_tech}"
    pH.font.bold = True
    pH.font.size = Pt(32)
    pH.font.color.rgb = RGBColor(255, 255, 255)
    pH.alignment = PP_ALIGN.CENTER

    # "НЕСЕ ЕМЕС?" (VS) бөлігі (Қызыл/Сұр)
    txNot = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 1.4), Inches(5.2), Inches(0.5))
    tNot = txNot.text_frame
    pNot = tNot.paragraphs[0]
    pNot.text = f"❌ Неге {not_tech} емес?"
    pNot.font.bold = True
    pNot.font.size = Pt(16)
    pNot.font.color.rgb = RGBColor(192, 57, 43)

    # "НЕГЕ ТАҢДАЛДЫ?" бөлігі
    txWhy = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 2.0), Inches(5.2), Inches(2.2))
    tWhy = txWhy.text_frame
    tWhy.word_wrap = True
    
    pWT = tWhy.paragraphs[0]
    pWT.text = f"✅ {reason_title}"
    pWT.font.bold = True
    pWT.font.size = Pt(18)
    pWT.font.color.rgb = RGBColor(*accent_color)
    
    pWD = tWhy.add_paragraph()
    pWD.text = reason_text
    pWD.font.size = Pt(15)
    pWD.font.color.rgb = RGBColor(52, 73, 94)
    pWD.space_before = Pt(8)

# --- PYTHON CARD ---
create_vs_card(
    left=0.6, top=1.4,
    icon="🐍", main_tech="PYTHON", not_tech="Java немесе C#",
    reason_title="AI мен ML экожүйесі",
    reason_text="Жобамыздың \"миы\" — Жасанды Интеллект (Scikit-learn, MediaPipe). Бұл салада Python — әлемдік стандарт. Java немесе C# қолдансақ, AI модельдерді қосу өте қиын болар еді. Python-да модельдерді серверге тікелей, еш кідіріссіз (native) орнаттық.",
    main_color=(41, 128, 185),    # Көк
    accent_color=(31, 97, 141)
)

# --- FASTAPI CARD ---
create_vs_card(
    left=6.9, top=1.4,
    icon="⚡", main_tech="FASTAPI", not_tech="Django немесе Flask",
    reason_title="Асинхронды жылдамдық (ASGI)",
    reason_text="Django — үлкен, монолитті әрі ауыр. Бізге камерадан секундына келетін 30 кадр (FPS) сұраныстарды қатырып тастамай өңдейтін өте жылдам сервер керек болды. FastAPI NodeJS деңгейіндегі жылдамдықты және автоматты Swagger (OpenAPI) құжаттамасын берді.",
    main_color=(39, 174, 96),     # Жасыл
    accent_color=(30, 132, 73)
)

# ============================================================
# 3. ТӨМЕНДЕГІ ФОРМУЛА (Қорытынды)
# ============================================================
formula_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(6.1), Inches(12.1), Inches(1.0)
)
formula_bg.fill.solid()
formula_bg.fill.fore_color.rgb = RGBColor(235, 222, 240)
formula_bg.line.color.rgb = RGBColor(142, 68, 173)
formula_bg.line.width = Pt(2)

txFoot = slide.shapes.add_textbox(Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.8))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "Python (AI үшін)  +  FastAPI (Асинхронды жылдамдық)  =  AI жобаларына арналған идеал Backend."
pFoot.font.bold = True
pFoot.font.size = Pt(20)
pFoot.font.color.rgb = RGBColor(118, 68, 138)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
