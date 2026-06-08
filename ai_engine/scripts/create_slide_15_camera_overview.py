from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_15_camera_dataset_and_models.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(12.73), Inches(0.7))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Камера жаттығулары: Деректер жинау, модельдерді салыстыру, таңдау"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ЖОҒАРҒЫ — 3 датасет карточкасы (Squat, Push-up, Lunge)
# ============================================================
section_title_1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(0.4))
tf_s1 = section_title_1.text_frame
p_s1 = tf_s1.paragraphs[0]
p_s1.text = "1. Деректер жинау (Datasets)"
p_s1.font.bold = True
p_s1.font.size = Pt(20)
p_s1.font.color.rgb = RGBColor(41, 128, 185)

datasets = [
    {
        "name": "Squat",
        "videos": "36 видео",
        "rows": "3 248 кадр",
        "features": "276 белгі",
        "source": "Өзіміз түсірген + YouTube",
        "color_bg": (232, 248, 245),
        "color_border": (26, 188, 156),
    },
    {
        "name": "Push-up",
        "videos": "100 видео",
        "rows": "5 489 кадр",
        "features": "152 белгі",
        "source": "Kaggle + YouTube",
        "color_bg": (235, 245, 251),
        "color_border": (52, 152, 219),
    },
    {
        "name": "Lunge",
        "videos": "3 009 видео",
        "rows": "3 174 кадр",
        "features": "100 белгі",
        "source": "REHAB24 + Синтетикалық",
        "color_bg": (254, 249, 231),
        "color_border": (241, 196, 15),
    },
]

card_w = Inches(4.0)
card_h = Inches(1.8)
top = Inches(1.45)

for i, ds in enumerate(datasets):
    left = Inches(0.5 + i * 4.2)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, card_w, card_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*ds["color_bg"])
    shape.line.color.rgb = RGBColor(*ds["color_border"])
    shape.line.width = Pt(2)

    tx = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), card_w - Inches(0.3), card_h - Inches(0.2))
    t = tx.text_frame
    t.word_wrap = True

    p_name = t.paragraphs[0]
    p_name.text = ds["name"]
    p_name.font.bold = True
    p_name.font.size = Pt(20)
    p_name.font.color.rgb = RGBColor(44, 62, 80)
    p_name.alignment = PP_ALIGN.CENTER

    for label, value in [("Видео:", ds["videos"]), ("Кадрлар:", ds["rows"]),
                          ("Белгілер:", ds["features"]), ("Деректер көзі:", ds["source"])]:
        pp = t.add_paragraph()
        pp.text = f"{label} {value}"
        pp.font.size = Pt(13)
        pp.font.color.rgb = RGBColor(52, 73, 94)

# ============================================================
# 3. ОРТАДА — Модельдерді салыстыру кестесі
# ============================================================
section_title_2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(8), Inches(0.4))
tf_s2 = section_title_2.text_frame
p_s2 = tf_s2.paragraphs[0]
p_s2.text = "2. Модельдерді салыстыру нәтижесі (Model Comparison)"
p_s2.font.bold = True
p_s2.font.size = Pt(20)
p_s2.font.color.rgb = RGBColor(41, 128, 185)

# Кесте
from pptx.util import Emu
rows_data = [
    ("Модель / Жаттығу", "Squat", "Push-up", "Lunge"),
    ("Random Forest", "85.8% ✅", "95.4% ✅", "96.6%"),
    ("Extra Trees", "84.2%", "92.2%", "95.4%"),
    ("Gradient Boosting", "83.9%", "94.6%", "97.3% ✅"),
]

table_shape = slide.shapes.add_table(
    rows=len(rows_data), cols=4,
    left=Inches(0.5), top=Inches(3.85),
    width=Inches(8.0), height=Inches(2.0)
)
table = table_shape.table

# Кесте стилі
for row_idx, row_data in enumerate(rows_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(15)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if row_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif "✅" in cell_text:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(39, 174, 96)
            else:
                paragraph.font.color.rgb = RGBColor(52, 73, 94)
        
        # Header row color
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

# ============================================================
# 4. ОҢ ЖАҚ — Неге осы алгоритмдер?
# ============================================================
why_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.0), Inches(3.85), Inches(4.0), Inches(3.0)
)
why_shape.fill.solid()
why_shape.fill.fore_color.rgb = RGBColor(244, 236, 247)
why_shape.line.color.rgb = RGBColor(142, 68, 173)
why_shape.line.width = Pt(2)

txWhy = slide.shapes.add_textbox(Inches(9.2), Inches(3.95), Inches(3.6), Inches(2.8))
tfW = txWhy.text_frame
tfW.word_wrap = True

pW_title = tfW.paragraphs[0]
pW_title.text = "Неге дәл осы алгоритмдер?"
pW_title.font.bold = True
pW_title.font.size = Pt(18)
pW_title.font.color.rgb = RGBColor(142, 68, 173)
pW_title.alignment = PP_ALIGN.CENTER

reasons = [
    "• Ансамбль әдістері (RF, GB) — жеке ағаштар шуылға (noise) тұрақты",
    "• Биомеханикалық белгілерде (бұрыштар, қашықтықтар) ең жақсы нәтиже береді",
    "• Нейрожеліге қарағанда жылдам жұмыс істейді (real-time)",
    "• Кішкене датасетте overfitting-ке аз ұшырайды",
]

for reason in reasons:
    pr = tfW.add_paragraph()
    pr.text = reason
    pr.font.size = Pt(13)
    pr.font.color.rgb = RGBColor(60, 40, 80)
    pr.space_before = Pt(6)

# ============================================================
# 5. ҚОРЫТЫНДЫ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4))
tfFoot = txFoot.text_frame
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Әр жаттығу үшін 3 алгоритмді салыстырып, дәлдігі ең жоғарысын таңдадық. "
    "Кестедегі ✅ белгісі — таңдалған модельді көрсетеді."
)
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(15)
pFoot.font.color.rgb = RGBColor(127, 140, 141)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
