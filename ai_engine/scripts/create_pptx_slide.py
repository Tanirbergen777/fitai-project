from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "biomechanics_slide.pptx"

prs = Presentation()
# 16:9 ratio
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

def add_card(slide, left, top, width, height, title, text, bg_rgb, border_rgb):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    shape.line.color.rgb = RGBColor(*border_rgb)
    shape.line.width = Pt(2)
    
    # Title
    txBox_title = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf_title = txBox_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.bold = True
    p_title.font.size = Pt(20)
    p_title.font.color.rgb = RGBColor(44, 62, 80)
    p_title.alignment = PP_ALIGN.CENTER
    
    # Body text
    txBox_text = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(1.0))
    tf_text = txBox_text.text_frame
    tf_text.word_wrap = True
    
    p_text = tf_text.paragraphs[0]
    p_text.text = text
    p_text.font.size = Pt(18)
    p_text.font.color.rgb = RGBColor(52, 73, 94)

# Data
cards = [
    {
        "title": "1. Squat (Отырып тұру)",
        "text": "• Тізе бұрышы (Knee Angle):\nЖамбас, тізе және тобық нүктелері арқылы есептеледі (90 градусқа дейін жетуі керек).\n\n• Арқаның түзулігі:\nИық пен жамбас арасындағы көлбеулік (алға еңкейіп кетпеуін бақылаймыз).",
        "bg": (232, 248, 245),
        "border": (26, 188, 156)
    },
    {
        "title": "2. Push-up (Жерден қисаю)",
        "text": "• Шынтақ бұрышы:\nИық, шынтақ және білезік нүктелерінен есептеледі (денені қаншалықты төмен түсіргенін бақылайды).\n\n• Дененің бір сызықта болуы:\nИық, жамбас және өкше нүктелерінің бір түзуде (180 градус) жатуын тексереді.",
        "bg": (235, 245, 251),
        "border": (52, 152, 219)
    },
    {
        "title": "3. Lunge (Алға қадам басу)",
        "text": "• Екі тізенің бұрышы:\nАлдыңғы тізе 90 градус болуы, ал артқы тізе жерге тимей тұруы бақыланады.\n\n• Дене симметриясы:\nСол және оң жақ иық пен жамбастың ауытқуы мен балансы есептеледі.",
        "bg": (254, 249, 231),
        "border": (241, 196, 15)
    }
]

# Add Title to Slide
txBox_main = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(13.33), Inches(1))
tf_main = txBox_main.text_frame
p_main = tf_main.paragraphs[0]
p_main.text = "Камера арқылы биомеханикалық бұрыштарды есептеу"
p_main.font.bold = True
p_main.font.size = Pt(28)
p_main.font.color.rgb = RGBColor(44, 62, 80)
p_main.alignment = PP_ALIGN.CENTER

# Add 3 cards
card_width = Inches(4)
card_height = Inches(4.5)
top_margin = Inches(1.8)

add_card(slide, Inches(0.5), top_margin, card_width, card_height, cards[0]["title"], cards[0]["text"], cards[0]["bg"], cards[0]["border"])
add_card(slide, Inches(4.7), top_margin, card_width, card_height, cards[1]["title"], cards[1]["text"], cards[1]["bg"], cards[1]["border"])
add_card(slide, Inches(8.9), top_margin, card_width, card_height, cards[2]["title"], cards[2]["text"], cards[2]["bg"], cards[2]["border"])

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
