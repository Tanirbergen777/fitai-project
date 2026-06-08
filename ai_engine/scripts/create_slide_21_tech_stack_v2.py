from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_21_tech_stack_v2.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жобаның Технологиялық Стекі және Жалпы Схемасы"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. СОЛ ЖАҚ: ТЕХНОЛОГИЯЛАР ТІЗІМІ
# ============================================================
list_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.2), Inches(3.5), Inches(5.8)
)
list_shape.fill.solid()
list_shape.fill.fore_color.rgb = RGBColor(244, 246, 247)
list_shape.line.color.rgb = RGBColor(189, 195, 199)
list_shape.line.width = Pt(1.5)

txList = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(3.3), Inches(5.6))
tL = txList.text_frame
tL.word_wrap = True

pL_title = tL.paragraphs[0]
pL_title.text = "Қолданылған технологиялар"
pL_title.font.bold = True
pL_title.font.size = Pt(20)
pL_title.font.color.rgb = RGBColor(44, 62, 80)
pL_title.alignment = PP_ALIGN.CENTER

techs = [
    ("Frontend:", "React, TypeScript, Vite"),
    ("Backend:", "Python, FastAPI"),
    ("Database:", "Neon (PostgreSQL)"),
    ("Auth & Media:", "Supabase"),
    ("AI / ML:", "Groq (Llama 3.3), MediaPipe, Scikit-learn"),
    ("Hosting:", "Render")
]

for title, desc in techs:
    pt = tL.add_paragraph()
    pt.text = f"• {title}"
    pt.font.bold = True
    pt.font.size = Pt(16)
    pt.font.color.rgb = RGBColor(41, 128, 185)
    pt.space_before = Pt(12)
    
    pd = tL.add_paragraph()
    pd.text = f"  {desc}"
    pd.font.size = Pt(14)
    pd.font.color.rgb = RGBColor(52, 73, 94)

# ============================================================
# 3. ОҢ ЖАҚ: TREE ДИАГРАММА
# ============================================================
def create_card(left, top, width, height, title, subtitle, bg_color, border_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = Pt(3)

    tx = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(width), Inches(height))
    t = tx.text_frame
    t.word_wrap = True
    
    p1 = t.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(*border_color)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = t.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(52, 73, 94)
    p2.alignment = PP_ALIGN.CENTER
    
    return shape

def add_connector_with_text(x1, y1, x2, y2, text, text_y_offset=0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.width = Pt(2)
    connector.line.color.rgb = RGBColor(149, 165, 166)
    
    if text:
        tx = slide.shapes.add_textbox(Inches((x1+x2)/2 - 1.0), Inches((y1+y2)/2 - 0.2 + text_y_offset), Inches(2.0), Inches(0.4))
        t = tx.text_frame
        p = t.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)
        p.alignment = PP_ALIGN.CENTER

# 3.1 FRONTEND
front_x = 7.5
front_y = 1.2
front_w = 3.0
front_h = 1.0
create_card(
    front_x, front_y, front_w, front_h,
    "💻 Frontend", "(React / Vite / TS)\nUser Interface", 
    (235, 245, 251), (41, 128, 185)
)

# 3.2 RENDER CLOUD CONTAINER
render_x = 6.0
render_y = 3.0
render_w = 6.0
render_h = 2.0
render_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(render_x), Inches(render_y), Inches(render_w), Inches(render_h)
)
render_bg.fill.solid()
render_bg.fill.fore_color.rgb = RGBColor(244, 246, 247)
render_bg.line.color.rgb = RGBColor(149, 165, 166)
render_bg.line.width = Pt(1)

txR = slide.shapes.add_textbox(Inches(render_x), Inches(render_y), Inches(render_w), Inches(0.4))
tR = txR.text_frame
pR = tR.paragraphs[0]
pR.text = "☁️ Hosted on Render Cloud"
pR.font.size = Pt(12)
pR.font.bold = True
pR.font.color.rgb = RGBColor(127, 140, 141)

# 3.3 BACKEND
back_x = 7.5
back_y = 3.5
back_w = 3.0
back_h = 1.0
create_card(
    back_x, back_y, back_w, back_h,
    "⚙️ Backend", "(FastAPI / Python)\nAPI Gateway & Core Logic", 
    (232, 248, 245), (39, 174, 96)
)

# Connect Front -> Back
add_connector_with_text(9.0, 2.2, 9.0, 3.5, "REST API / HTTP\n(Сұраныстар)")

# 3.4 DATABASE (LEFT)
db_x = 4.5
db_y = 6.0
db_w = 2.8
db_h = 1.0
create_card(
    db_x, db_y, db_w, db_h,
    "🗄️ Database", "(Neon / PostgreSQL)\nUser & Plan Data", 
    (253, 235, 208), (211, 84, 0)
)
# Connect Back -> DB
connector1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(8.5), Inches(4.5), Inches(5.9), Inches(6.0))
connector1.line.width = Pt(2)
connector1.line.color.rgb = RGBColor(149, 165, 166)
add_connector_with_text(5.9, 5.0, 5.9, 5.0, "SQL Queries\n(Деректер сақтау)", -0.4)


# 3.5 SUPABASE (CENTER)
supa_x = 7.6
supa_y = 6.0
supa_w = 2.8
supa_h = 1.0
create_card(
    supa_x, supa_y, supa_w, supa_h,
    "🔒 Auth & Storage", "(Supabase)\nAuthentication & Media", 
    (232, 218, 239), (142, 68, 173)
)
# Connect Back -> Supabase
connector2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.0), Inches(4.5), Inches(9.0), Inches(6.0))
connector2.line.width = Pt(2)
connector2.line.color.rgb = RGBColor(149, 165, 166)
add_connector_with_text(9.0, 5.25, 9.0, 5.25, "API / Bucket\n(Токендер & Файлдар)")


# 3.6 AI ENGINE (RIGHT)
ai_x = 10.7
ai_y = 6.0
ai_w = 2.4
ai_h = 1.0
create_card(
    ai_x, ai_y, ai_w, ai_h,
    "🧠 AI Engine", "(Groq / MediaPipe)\nML Models", 
    (253, 237, 236), (192, 57, 43)
)
# Connect Back -> AI
connector3 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(9.5), Inches(4.5), Inches(11.9), Inches(6.0))
connector3.line.width = Pt(2)
connector3.line.color.rgb = RGBColor(149, 165, 166)
add_connector_with_text(11.9, 5.0, 11.9, 5.0, "Inference\n(Жүйке желісі)", -0.4)

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
