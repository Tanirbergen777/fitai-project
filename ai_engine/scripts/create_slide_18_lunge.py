from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_18_lunge_model.pptx"
img_path = OUTPUT_DIR / "lunge_cm_presentation.png"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Lunge (Алға қадам басу) — Деректерді ұлғайту және нәтижесі"
p.font.bold = True
p.font.size = Pt(30)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ЖОҒАРҒЫ СИПАТТАМА
# ============================================================
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.73), Inches(0.7))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = (
    "Бастапқыда деректер аз болып (174 қатар), дәлдік 60% ғана болды. "
    "Биомеханикалық зерттеулерге сүйеніп синтетикалық деректер генерацияладық "
    "және деректер қорын 3 174 қатарға дейін ұлғайтып, модельді қайта оқыттық."
)
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(52, 73, 94)

# ============================================================
# 3. СОЛ ЖАҚ — «БҰРЫН vs КЕЙІН» карточкалары
# ============================================================

# --- БҰРЫН (қызыл карточка) ---
before_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(2.0), Inches(3.8), Inches(2.2)
)
before_shape.fill.solid()
before_shape.fill.fore_color.rgb = RGBColor(253, 237, 236)  # Light red
before_shape.line.color.rgb = RGBColor(231, 76, 60)
before_shape.line.width = Pt(2)

txBefore = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(3.4), Inches(2.0))
tfB = txBefore.text_frame
tfB.word_wrap = True

pB_title = tfB.paragraphs[0]
pB_title.text = "❌  БҰРЫН"
pB_title.font.bold = True
pB_title.font.size = Pt(22)
pB_title.font.color.rgb = RGBColor(192, 57, 43)
pB_title.alignment = PP_ALIGN.CENTER

pB1 = tfB.add_paragraph()
pB1.text = "• Деректер: 174 қатар (9 видео)"
pB1.font.size = Pt(16)
pB1.font.color.rgb = RGBColor(100, 30, 22)
pB1.space_before = Pt(12)

pB2 = tfB.add_paragraph()
pB2.text = "• Алгоритм: RandomForest"
pB2.font.size = Pt(16)
pB2.font.color.rgb = RGBColor(100, 30, 22)

pB3 = tfB.add_paragraph()
pB3.text = "• Дәлдік: 60.0%"
pB3.font.size = Pt(20)
pB3.font.bold = True
pB3.font.color.rgb = RGBColor(192, 57, 43)

# --- КЕЙІН (жасыл карточка) ---
after_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(4.5), Inches(3.8), Inches(2.2)
)
after_shape.fill.solid()
after_shape.fill.fore_color.rgb = RGBColor(234, 250, 241)  # Light green
after_shape.line.color.rgb = RGBColor(39, 174, 96)
after_shape.line.width = Pt(2)

txAfter = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(3.4), Inches(2.0))
tfA = txAfter.text_frame
tfA.word_wrap = True

pA_title = tfA.paragraphs[0]
pA_title.text = "✅  КЕЙІН"
pA_title.font.bold = True
pA_title.font.size = Pt(22)
pA_title.font.color.rgb = RGBColor(39, 174, 96)
pA_title.alignment = PP_ALIGN.CENTER

pA1 = tfA.add_paragraph()
pA1.text = "• Деректер: 3 174 қатар (+3000 синтетикалық)"
pA1.font.size = Pt(16)
pA1.font.color.rgb = RGBColor(20, 90, 50)
pA1.space_before = Pt(12)

pA2 = tfA.add_paragraph()
pA2.text = "• Алгоритм: GradientBoosting"
pA2.font.size = Pt(16)
pA2.font.color.rgb = RGBColor(20, 90, 50)

pA3 = tfA.add_paragraph()
pA3.text = "• Дәлдік: 97.3% (+37.3%)"
pA3.font.size = Pt(20)
pA3.font.bold = True
pA3.font.color.rgb = RGBColor(39, 174, 96)

# --- Стрелка (БҰРЫН → КЕЙІН) ---
arrow = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW,
    Inches(2.1), Inches(4.2), Inches(0.6), Inches(0.3)
)
arrow.fill.solid()
arrow.fill.fore_color.rgb = RGBColor(41, 128, 185)
arrow.line.fill.background()

# ============================================================
# 4. ОРТАДА — Confusion Matrix суреті
# ============================================================
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(4.8), Inches(2.2), height=Inches(4.2))

# ============================================================
# 5. ОҢ ЖАҚ — Қате типтері
# ============================================================
error_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.7)
)
error_shape.fill.solid()
error_shape.fill.fore_color.rgb = RGBColor(235, 245, 251)  # Light blue
error_shape.line.color.rgb = RGBColor(41, 128, 185)
error_shape.line.width = Pt(2)

txErr = slide.shapes.add_textbox(Inches(9.4), Inches(2.1), Inches(3.4), Inches(4.5))
tfE = txErr.text_frame
tfE.word_wrap = True

pE_title = tfE.paragraphs[0]
pE_title.text = "4 типтегі қателік:"
pE_title.font.bold = True
pE_title.font.size = Pt(20)
pE_title.font.color.rgb = RGBColor(41, 128, 185)
pE_title.alignment = PP_ALIGN.CENTER

errors = [
    ("1. Knee Over Toe", "Тізе аяқ ұшынан алға шығып кетуі"),
    ("2. Shallow Lunge", "Тізені жеткілікті бүкпей таяз отыру"),
    ("3. Excessive Lean", "Дененің алға тым еңкейіп кетуі"),
    ("4. Knee Valgus", "Тізенің ішке қарай ауытқуы"),
]

for title, desc in errors:
    pT = tfE.add_paragraph()
    pT.text = title
    pT.font.bold = True
    pT.font.size = Pt(16)
    pT.font.color.rgb = RGBColor(44, 62, 80)
    pT.space_before = Pt(10)
    
    pD = tfE.add_paragraph()
    pD.text = desc
    pD.font.size = Pt(14)
    pD.font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# 6. ҚОРЫТЫНДЫ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.5))
tfFoot = txFoot.text_frame
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Қорытынды: Деректерді биомеханикалық білім негізінде ұлғайту (Data Augmentation) "
    "арқылы модельдің дәлдігін 60%-дан 97.3%-ға дейін көтеруге қол жеткіздік."
)
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(16)
pFoot.font.color.rgb = RGBColor(142, 68, 173)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
