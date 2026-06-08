from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_20_ai_chat.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "AI Chat: Неге сыртқы LLM (Groq + Llama 3.3 70B) таңдалды?"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ОРТАДА — 3 СЕБЕП (Карточкалар)
# ============================================================
cards = [
    {
        "icon": "⚡",
        "title": "Жылдамдық (Ultra-low Latency)",
        "desc": "Біз OpenAI (ChatGPT) емес, Groq платформасын таңдадық. Себебі Groq дәстүрлі GPU емес, арнайы LPU (Language Processing Unit) процессорларын қолданады. Бұл чаттың пайдаланушыға секундына 800+ токен жылдамдықпен (лезде) жауап беруін қамтамасыз етеді.",
        "bg_color": (254, 249, 231),
        "line_color": (241, 196, 15)
    },
    {
        "icon": "💻",
        "title": "Ресурс үнемдеу (Hardware)",
        "desc": "Llama 3.3 (70 Billion parameters) сияқты алып модельді өз серверімізде іске қосу үшін кем дегенде 4-8 қымбат видеокарта (A100) қажет болар еді. Бұл өте тиімсіз. Сыртқы API қолдану серверлік шығындарды 99%-ға азайтты.",
        "bg_color": (235, 245, 251),
        "line_color": (52, 152, 219)
    },
    {
        "icon": "🧠",
        "title": "Сапа және Контекст",
        "desc": "Llama 3.3 70B — қазіргі таңдағы ең мықты ашық (open-source) модельдердің бірі. Ол фитнес, диетология терминдерін жақсы түсінеді және пайдаланушыға тек ақпарат беріп қана қоймай, эмпатиямен (мотивация беріп) сөйлесе алады.",
        "bg_color": (232, 248, 245),
        "line_color": (39, 174, 96)
    }
]

card_w = Inches(3.8)
card_h = Inches(4.0)
top = Inches(1.5)

for i, card in enumerate(cards):
    left = Inches(0.6 + i * 4.1)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, card_w, card_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*card["bg_color"])
    shape.line.color.rgb = RGBColor(*card["line_color"])
    shape.line.width = Pt(2)

    tx = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.4))
    t = tx.text_frame
    t.word_wrap = True

    p_icon = t.paragraphs[0]
    p_icon.text = f"{card['icon']} {card['title']}"
    p_icon.font.bold = True
    p_icon.font.size = Pt(20)
    p_icon.font.color.rgb = RGBColor(*card["line_color"])
    p_icon.alignment = PP_ALIGN.CENTER
    
    p_space = t.add_paragraph()
    p_space.text = ""
    p_space.font.size = Pt(10)

    p_desc = t.add_paragraph()
    p_desc.text = card["desc"]
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = RGBColor(52, 73, 94)

# ============================================================
# 3. ТӨМЕНДЕ — ҚОРЫТЫНДЫ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.33), Inches(1.0))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Шешім: Groq API + Llama 3.3 70B = Серверге салмақ түсірмейтін, "
    "максималды жылдам әрі ақылды чат архитектурасы."
)
pFoot.font.bold = True
pFoot.font.size = Pt(20)
pFoot.font.color.rgb = RGBColor(142, 68, 173)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
