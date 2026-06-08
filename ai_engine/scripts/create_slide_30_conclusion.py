from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_30_conclusion.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.0))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Қорытынды: Мақсатқа толықтай қол жеткізілді"
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. НӘТИЖЕЛЕР (Checklist)
# ============================================================
def create_check_card(y, title, desc):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(y), Inches(10.33), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 249, 249)
    card.line.color.rgb = RGBColor(39, 174, 96) # Green
    card.line.width = Pt(2)
    
    # Checkmark icon
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.7), Inches(y + 0.2), Inches(0.8), Inches(0.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(39, 174, 96)
    circle.line.fill.background()
    
    tC = circle.text_frame.paragraphs[0]
    tC.text = "✓"
    tC.font.bold = True
    tC.font.size = Pt(36)
    tC.font.color.rgb = RGBColor(255, 255, 255)
    tC.alignment = PP_ALIGN.CENTER
    
    # Text content
    txT = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.1), Inches(9.0), Inches(1.0))
    tT = txT.text_frame
    tT.word_wrap = True
    
    pT = tT.paragraphs[0]
    pT.text = title
    pT.font.bold = True
    pT.font.size = Pt(20)
    pT.font.color.rgb = RGBColor(44, 62, 80)
    
    pD = tT.add_paragraph()
    pD.text = desc
    pD.font.size = Pt(16)
    pD.font.color.rgb = RGBColor(127, 140, 141)

# Item 1
create_check_card(1.8, "Жобаның негізгі мақсаты орындалды", 
                  "Жасанды интеллект (AI) пен Компьютерлік көру (CV) технологияларына негізделген толыққанды, жекелендірілген фитнес экожүйесі құрылды.")

# Item 2
create_check_card(3.2, "Технологиялық интеграция сәтті аяқталды", 
                  "4 түрлі AI моделі (MediaPipe, ML, LLM) Backend (FastAPI) және Frontend (React) бөліктерімен мінсіз байланыстырылды.")

# Item 3
create_check_card(4.6, "Әлеуметтік маңызы бар мәселе шешілді", 
                  "Спорт қолжетімді, ал жаттығу қауіпсіз етілді. Жеке бапкердің орнын басатын инновациялық цифрлық шешім нарыққа ұсынылды.")

# ============================================================
# 3. ФИНАЛ (Қоштасу)
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.33), Inches(1.0))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "Назарларыңызға рақмет!\nСұрақтарыңызға жауап беруге дайынмын."
pFoot.font.bold = True
pFoot.font.size = Pt(28)
pFoot.font.color.rgb = RGBColor(41, 128, 185) # Blue
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
