from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "squat_model_slide.pptx"
img_path = OUTPUT_DIR / "squat_cm_presentation.png"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 1. Title
txBox_title = slide.shapes.add_textbox(Inches(0), Inches(0.4), Inches(13.33), Inches(1))
tf_title = txBox_title.text_frame
p_title = tf_title.paragraphs[0]
p_title.text = "Squat (Отырып тұру) моделін оқыту және нәтижесі"
p_title.font.bold = True
p_title.font.size = Pt(32)
p_title.font.color.rgb = RGBColor(44, 62, 80)
p_title.alignment = PP_ALIGN.CENTER

# 2. Top description
txBox_desc = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.33), Inches(1))
tf_desc = txBox_desc.text_frame
tf_desc.word_wrap = True
p_desc = tf_desc.paragraphs[0]
p_desc.text = "Модельді оқыту үшін 3 248 кадрдан (row) тұратын арнайы датасет жиналды. Оның ішінде 'Дұрыс' жасалған және 'Арқаны бүгіп/Тізені қате бүгіп' жасалған кадрлар тең дәрежеде (balanced) қамтылды."
p_desc.font.size = Pt(18)
p_desc.font.color.rgb = RGBColor(52, 73, 94)

# 3. Add Image
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(2.6), height=Inches(3.8))

# 4. Results text box
txBox_res = slide.shapes.add_textbox(Inches(6.5), Inches(2.6), Inches(6), Inches(3.8))
tf_res = txBox_res.text_frame
tf_res.word_wrap = True

def add_bullet(tf, bold_text, normal_text):
    p = tf.add_paragraph()
    p.text = bold_text
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(41, 128, 185)
    
    run = p.add_run()
    run.text = normal_text
    run.font.bold = False
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(52, 73, 94)
    p.space_after = Pt(14)

tf_res.text = "" # clear default paragraph
add_bullet(tf_res, "• Алгоритм: ", "Random Forest")
add_bullet(tf_res, "• Дәлдік (Accuracy): ", "85.8%")
add_bullet(tf_res, "• Feature Count: ", "276 түрлі биомеханикалық белгі (нүктелер, координаталар, бұрыштар).")
add_bullet(tf_res, "• Талдау: ", "Модель қате жасалған squat-ты өте жақсы табады (Матрицада 'Wrong' класын табу деңгейі жоғары).")

# 5. Conclusion (Footer)
txBox_foot = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.8))
tf_foot = txBox_foot.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.text = "Қорытынды: Squat техникасы күрделі болғанымен, гибридті модель (MediaPipe + RF) 85.8% дәлдікпен сенімді жұмыс істейді."
p_foot.font.bold = True
p_foot.font.italic = True
p_foot.font.size = Pt(18)
p_foot.font.color.rgb = RGBColor(39, 174, 96) # Green
p_foot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
