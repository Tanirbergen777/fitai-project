from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_26_results_metrics.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жалпы нәтижелер: Барлық метрика бір жерде"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. МЕТРИКАЛАР КАРТОЧКАЛАРЫ (Dashboard)
# ============================================================
def create_metric_card(x, y, w, h, icon, number, title, desc, main_color):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 249, 249)
    card.line.color.rgb = RGBColor(*main_color)
    card.line.width = Pt(2)
    
    # Icon and Number
    txN = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(w), Inches(0.8))
    tN = txN.text_frame
    pN = tN.paragraphs[0]
    pN.text = f"{icon} {number}"
    pN.font.bold = True
    pN.font.size = Pt(40)
    pN.font.color.rgb = RGBColor(*main_color)
    pN.alignment = PP_ALIGN.CENTER
    
    # Title
    txT = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(w), Inches(0.5))
    pT = txT.text_frame.paragraphs[0]
    pT.text = title
    pT.font.bold = True
    pT.font.size = Pt(20)
    pT.font.color.rgb = RGBColor(44, 62, 80)
    pT.alignment = PP_ALIGN.CENTER
    
    # Description
    txD = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.6), Inches(w - 0.4), Inches(h - 1.6))
    tD = txD.text_frame
    tD.word_wrap = True
    pD = tD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(14)
    pD.font.color.rgb = RGBColor(127, 140, 141)
    pD.alignment = PP_ALIGN.CENTER

# Define Layout for 4 cards (2x2 grid)
w, h = 5.6, 2.5
x1, x2 = 0.8, 6.9
y1, y2 = 1.5, 4.3

# Card 1: Accuracy
create_metric_card(
    x1, y1, w, h, 
    "🎯", "97%+", "Орташа Дәлдік (Accuracy)", 
    "Компьютерлік көру модельдері (Lunge 97.3%, Push-up 95.4%, Squat 96%+) және фитнес алгоритмдері.",
    (39, 174, 96) # Green
)

# Card 2: Datasets (Updated based on user feedback)
create_metric_card(
    x2, y1, w, h, 
    "📊", "20 000+", "Жиналған Датасеттер", 
    "CV үшін 11 900+ кадр. Ұсыныс жүйесіне (Diet/Workout Recommender) арналған NHANES (8 790+) деректері.",
    (41, 128, 185) # Blue
)

# Card 3: Backend
create_metric_card(
    x1, y2, w, h, 
    "⚙️", "30+", "API Endpoints (FastAPI)", 
    "Жүйе микросервис ретінде жұмыс істейтін 7 бөлек Router-ге және қатаң SQL байланыстарға бөлінген.",
    (211, 84, 0) # Orange
)

# Card 4: AI Tech
create_metric_card(
    x2, y2, w, h, 
    "🤖", "4 Түрлі", "AI Интеграциясы", 
    "Жобада Machine Learning (Random Forest, Gradient Boosting), MediaPipe (3D Tracking) және LLM (Llama 3.3) біріктірілді.",
    (142, 68, 173) # Purple
)

# ============================================================
# 3. ҚОРЫТЫНДЫ (Highlight)
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.5))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "Қорытынды: Бұл жай ғана прототип емес, өндіріске (Production) толық дайын Кешенді AI Фитнес Экожүйесі."
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(16)
pFoot.font.color.rgb = RGBColor(192, 57, 43)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
