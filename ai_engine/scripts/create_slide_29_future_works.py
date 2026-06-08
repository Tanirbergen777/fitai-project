from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_29_future_works.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жобаның болашағы: Келесі стратегиялық қадамдар"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ROADMAP КАРТОЧКАЛАРЫ
# ============================================================
def create_future_card(x, y, w, h, icon, title, desc, border_color, bg_color):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*bg_color)
    card.line.color.rgb = RGBColor(*border_color)
    card.line.width = Pt(2)
    
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + w/2 - 0.4), Inches(y - 0.4), Inches(0.8), Inches(0.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*border_color)
    circle.line.fill.background()
    
    tC = circle.text_frame.paragraphs[0]
    tC.text = icon
    tC.font.size = Pt(24)
    tC.alignment = PP_ALIGN.CENTER
    
    # Title
    txT = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(w), Inches(0.5))
    pT = txT.text_frame.paragraphs[0]
    pT.text = title
    pT.font.bold = True
    pT.font.size = Pt(18)
    pT.font.color.rgb = RGBColor(44, 62, 80)
    pT.alignment = PP_ALIGN.CENTER
    
    # Description
    txD = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.0), Inches(w - 0.2), Inches(h - 1.1))
    tD = txD.text_frame
    tD.word_wrap = True
    pD = tD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(14)
    pD.font.color.rgb = RGBColor(52, 73, 94)
    pD.alignment = PP_ALIGN.CENTER

w, h = 2.8, 3.0
y_pos = 2.5
spacing = 3.2

# 1. Dataset Expansion
create_future_card(
    1.0, y_pos, w, h, "📊", "Датасет Кеңейту",
    "Компьютерлік көру модельдеріне Deadlift, Pull-up, Plank сияқты 20-дан астам жаңа жаттығу түрін қосу.",
    (41, 128, 185), (235, 245, 251) # Blue
)

# 2. Mobile App
create_future_card(
    1.0 + spacing*1, y_pos, w, h, "📱", "Мобильді Қосымша",
    "React Native арқылы жүйені толыққанды iOS және Android қосымшасына айналдырып, App Store-ға шығару.",
    (39, 174, 96), (232, 248, 245) # Green
)

# 3. Health Monitoring
create_future_card(
    1.0 + spacing*2, y_pos, w, h, "⌚", "Smart Мониторинг",
    "Apple Health және Google Fit-пен интеграция жасап, смарт-сағаттардан пульс пен калорияны автоматты оқу.",
    (211, 84, 0), (253, 235, 208) # Orange
)

# 4. AI Referee (NEW)
create_future_card(
    1.0 + spacing*3, y_pos, w, h, "⚖️", "AI Судья (Referee)",
    "Тек фитнес емес, кәсіби спорт түрлеріне (пауэрлифтинг, гимнастика) арналған әділ 'AI Судья' жүйесін жасау.",
    (142, 68, 173), (232, 218, 239) # Purple
)

# Draw connecting arrows between them
for i in range(3):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(1.0 + w + spacing*i + 0.1), Inches(y_pos + h/2 - 0.2), Inches(0.2), Inches(0.4)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(189, 195, 199)
    arrow.line.fill.background()

# ============================================================
# 3. ҚОРЫТЫНДЫ (Highlight)
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "FitAI — бұл аяқталған дипломдық жұмыс емес, нарыққа шығуға және инвестиция тартуға дайын үлкен IT-стартаптың іргетасы."
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(18)
pFoot.font.color.rgb = RGBColor(192, 57, 43)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
