from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_19_architecture.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Жүйе архитектурасы: AI модельдерін FastAPI арқылы интеграциялау"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. АРХИТЕКТУРА СХЕМАСЫ (3 блок)
# ============================================================

# Блок 1: Frontend
frontend_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.8), Inches(3.2), Inches(3.0)
)
frontend_shape.fill.solid()
frontend_shape.fill.fore_color.rgb = RGBColor(235, 245, 251)
frontend_shape.line.color.rgb = RGBColor(41, 128, 185)
frontend_shape.line.width = Pt(2)

txFront = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(3.0), Inches(2.8))
tfF = txFront.text_frame
tfF.word_wrap = True
pF_title = tfF.paragraphs[0]
pF_title.text = "📱 Frontend"
pF_title.font.bold = True
pF_title.font.size = Pt(22)
pF_title.font.color.rgb = RGBColor(41, 128, 185)
pF_title.alignment = PP_ALIGN.CENTER
tfF.add_paragraph()
pF1 = tfF.add_paragraph()
pF1.text = "• User Input (Деректер енгізу)"
pF1.font.size = Pt(16)
pF1.font.color.rgb = RGBColor(52, 73, 94)
pF2 = tfF.add_paragraph()
pF2.text = "• Camera Feed (Видео жіберу)"
pF2.font.size = Pt(16)
pF2.font.color.rgb = RGBColor(52, 73, 94)
pF3 = tfF.add_paragraph()
pF3.text = "• Нәтижелерді визуалдау"
pF3.font.size = Pt(16)
pF3.font.color.rgb = RGBColor(52, 73, 94)

# Arrow 1
arrow1 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW,
    Inches(4.1), Inches(3.1), Inches(0.8), Inches(0.4)
)
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(189, 195, 199)
arrow1.line.fill.background()

# Блок 2: Backend
backend_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.0), Inches(1.8), Inches(3.2), Inches(3.0)
)
backend_shape.fill.solid()
backend_shape.fill.fore_color.rgb = RGBColor(232, 248, 245)
backend_shape.line.color.rgb = RGBColor(39, 174, 96)
backend_shape.line.width = Pt(2)

txBack = slide.shapes.add_textbox(Inches(5.1), Inches(1.9), Inches(3.0), Inches(2.8))
tfB = txBack.text_frame
tfB.word_wrap = True
pB_title = tfB.paragraphs[0]
pB_title.text = "⚙️ Backend (FastAPI)"
pB_title.font.bold = True
pB_title.font.size = Pt(22)
pB_title.font.color.rgb = RGBColor(39, 174, 96)
pB_title.alignment = PP_ALIGN.CENTER
tfB.add_paragraph()
pB1 = tfB.add_paragraph()
pB1.text = "• /nutrition_plan"
pB1.font.size = Pt(16)
pB1.font.color.rgb = RGBColor(52, 73, 94)
pB2 = tfB.add_paragraph()
pB2.text = "• /workout_plan"
pB2.font.size = Pt(16)
pB2.font.color.rgb = RGBColor(52, 73, 94)
pB3 = tfB.add_paragraph()
pB3.text = "• /pose_estimation"
pB3.font.size = Pt(16)
pB3.font.color.rgb = RGBColor(52, 73, 94)

# Arrow 2
arrow2 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW,
    Inches(8.3), Inches(3.1), Inches(0.8), Inches(0.4)
)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(189, 195, 199)
arrow2.line.fill.background()

# Блок 3: AI Engine
ai_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.2), Inches(1.8), Inches(3.2), Inches(3.0)
)
ai_shape.fill.solid()
ai_shape.fill.fore_color.rgb = RGBColor(254, 249, 231)
ai_shape.line.color.rgb = RGBColor(241, 196, 15)
ai_shape.line.width = Pt(2)

txAI = slide.shapes.add_textbox(Inches(9.3), Inches(1.9), Inches(3.0), Inches(2.8))
tfAI = txAI.text_frame
tfAI.word_wrap = True
pAI_title = tfAI.paragraphs[0]
pAI_title.text = "🧠 AI Engine (.pkl)"
pAI_title.font.bold = True
pAI_title.font.size = Pt(22)
pAI_title.font.color.rgb = RGBColor(211, 84, 0)
pAI_title.alignment = PP_ALIGN.CENTER
tfAI.add_paragraph()
pAI1 = tfAI.add_paragraph()
pAI1.text = "• Random Forest / XGBoost"
pAI1.font.size = Pt(16)
pAI1.font.color.rgb = RGBColor(52, 73, 94)
pAI2 = tfAI.add_paragraph()
pAI2.text = "• MediaPipe (Keypoints)"
pAI2.font.size = Pt(16)
pAI2.font.color.rgb = RGBColor(52, 73, 94)
pAI3 = tfAI.add_paragraph()
pAI3.text = "• Feature Extraction"
pAI3.font.size = Pt(16)
pAI3.font.color.rgb = RGBColor(52, 73, 94)


# ============================================================
# 3. ТӨМЕНДЕ — Неге FastAPI?
# ============================================================
why_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.2), Inches(11.6), Inches(1.8)
)
why_shape.fill.solid()
why_shape.fill.fore_color.rgb = RGBColor(244, 236, 247)
why_shape.line.color.rgb = RGBColor(142, 68, 173)
why_shape.line.width = Pt(2)

txWhy = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.2), Inches(1.6))
tfWhy = txWhy.text_frame
tfWhy.word_wrap = True

pW_title = tfWhy.paragraphs[0]
pW_title.text = "Неге бұл сәулетте FastAPI таңдалды?"
pW_title.font.bold = True
pW_title.font.size = Pt(20)
pW_title.font.color.rgb = RGBColor(142, 68, 173)

reasons = [
    "• Жоғары жылдамдық: Нақты уақыттағы (real-time, 30 FPS) камера анализі үшін өте маңызды.",
    "• Асинхронды өңдеу (Asynchronous): Сервер бір уақытта келетін көптеген сұраныстардан құлап қалмайды.",
    "• Микросервис архитектурасы: AI модельдері (Machine Learning) мен Backend бір-бірімен өте жеңіл әрі сенімді байланысқан."
]

for r in reasons:
    pr = tfWhy.add_paragraph()
    pr.text = r
    pr.font.size = Pt(16)
    pr.font.color.rgb = RGBColor(52, 73, 94)
    pr.space_before = Pt(6)

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
