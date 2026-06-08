from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_demo_scenario.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Demo сценарий: Қолданушы жолы (User Flow)"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ПРОЦЕСС ҚАДАМДАРЫ (Timeline / Flowchart)
# ============================================================
steps = [
    {"icon": "🔐", "title": "Кіру", "desc": "Жүйеге тіркелу және қауіпсіз авторизация (JWT)"},
    {"icon": "👤", "title": "Профиль", "desc": "Бойы, салмағы, жасы және мақсатты енгізу"},
    {"icon": "🧠", "title": "AI Анализ", "desc": "Машиналық оқыту арқылы деректерді талдау"},
    {"icon": "🍎", "title": "Тамақтану", "desc": "AI жасаған жеке рацион және калория"},
    {"icon": "📷", "title": "Жаттығу", "desc": "Камера (CV) арқылы қателерді тікелей түзету"},
    {"icon": "🤖", "title": "AI Чат", "desc": "Фитнес-кеңесшімен (Llama 3.3) сұрақ-жауап"}
]

# Layout parameters for 2 rows of 3 items
start_x = 1.0
start_y = 2.0
x_step = 4.0
y_step = 2.3

# Colors
bg_colors = [(235, 245, 251), (253, 235, 208), (232, 248, 245), (253, 237, 236), (232, 218, 239), (252, 243, 207)]
border_colors = [(41, 128, 185), (211, 84, 0), (39, 174, 96), (192, 57, 43), (142, 68, 173), (241, 196, 15)]

def draw_step(index, x, y, step_data, bg_c, border_c):
    # Number circle
    num_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - 0.2), Inches(y - 0.2), Inches(0.6), Inches(0.6)
    )
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = RGBColor(*border_c)
    num_shape.line.fill.background()
    
    tNum = num_shape.text_frame
    pN = tNum.paragraphs[0]
    pN.text = str(index + 1)
    pN.font.bold = True
    pN.font.size = Pt(16)
    pN.alignment = PP_ALIGN.CENTER
    
    # Card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(3.2), Inches(1.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*bg_c)
    card.line.color.rgb = RGBColor(*border_c)
    card.line.width = Pt(2)
    
    # Text
    tx = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(3.0), Inches(1.3))
    t = tx.text_frame
    t.word_wrap = True
    
    p1 = t.paragraphs[0]
    p1.text = f"{step_data['icon']} {step_data['title']}"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = RGBColor(44, 62, 80)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = t.add_paragraph()
    p2.text = step_data['desc']
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(52, 73, 94)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # Bring number to front implicitly by drawing it slightly differently or leaving it
    # Actually, drawing it after the card ensures it's on top. Let's redraw number to be on top.
    num_shape2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - 0.2), Inches(y - 0.2), Inches(0.5), Inches(0.5)
    )
    num_shape2.fill.solid()
    num_shape2.fill.fore_color.rgb = RGBColor(*border_c)
    num_shape2.line.fill.background()
    tNum2 = num_shape2.text_frame
    pN2 = tNum2.paragraphs[0]
    pN2.text = str(index + 1)
    pN2.font.bold = True
    pN2.font.size = Pt(16)
    pN2.alignment = PP_ALIGN.CENTER

    return (x + 3.2, y + 0.75), (x + 1.6, y + 1.5) # right edge, bottom edge

def draw_arrow(pt1, pt2, style=MSO_CONNECTOR.STRAIGHT):
    conn = slide.shapes.add_connector(
        style, 
        Inches(pt1[0]), Inches(pt1[1]), 
        Inches(pt2[0]), Inches(pt2[1])
    )
    conn.line.width = Pt(4)
    conn.line.color.rgb = RGBColor(149, 165, 166)
    # Add arrowhead (end style)
    # python-pptx doesn't have direct arrowhead property easily exposed without XML manipulation, 
    # but the logic implies flow. We'll use simple connectors.

# Top Row: 1 -> 2 -> 3
edges = []
for i in range(3):
    r_edge, b_edge = draw_step(i, start_x + i * x_step, start_y, steps[i], bg_colors[i], border_colors[i])
    edges.append((r_edge, b_edge))

# Draw horizontal arrows for row 1
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + 3.3), Inches(start_y + 0.6), Inches(0.5), Inches(0.3)).fill.solid()
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + x_step + 3.3), Inches(start_y + 0.6), Inches(0.5), Inches(0.3)).fill.solid()

# Draw arrow from 3 to 4 (Down and left) - U-turn
# We can just draw a down arrow from 3 and right arrow for the bottom row, or read left-to-right again.
# Let's read left to right again for simplicity:
# Bottom Row: 4 -> 5 -> 6
for i in range(3, 6):
    r_edge, b_edge = draw_step(i, start_x + (i-3) * x_step, start_y + y_step, steps[i], bg_colors[i], border_colors[i])

# Draw horizontal arrows for row 2
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + 3.3), Inches(start_y + y_step + 0.6), Inches(0.5), Inches(0.3)).fill.solid()
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + x_step + 3.3), Inches(start_y + y_step + 0.6), Inches(0.5), Inches(0.3)).fill.solid()

# Connect Row 1 to Row 2
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(start_x + 2 * x_step + 1.4), Inches(start_y + 1.6), Inches(0.4), Inches(0.6)).fill.solid()

# ============================================================
# 3. ҚОРЫТЫНДЫ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = "Барлық қадамдар толығымен автоматтандырылған. Қолданушы тек өз мақсатын көрсетеді, қалғанын AI жасайды."
pFoot.font.bold = True
pFoot.font.italic = True
pFoot.font.size = Pt(18)
pFoot.font.color.rgb = RGBColor(41, 128, 185)
pFoot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
