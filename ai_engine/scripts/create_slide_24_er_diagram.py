from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "slide_24_er_diagram.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# 1. ТАҚЫРЫП
# ============================================================
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Мәліметтер базасы: Neon (PostgreSQL) және ER Диаграмма"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(44, 62, 80)
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 2. ER ДИАГРАММА (Кестелер)
# ============================================================

def draw_table(x, y, table_name, columns, header_color):
    width = 2.4
    height = 0.5 + len(columns) * 0.35
    
    # Main border
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 249, 249)
    bg.line.color.rgb = RGBColor(*header_color)
    bg.line.width = Pt(2)
    
    # Header block
    head = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(0.5)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(*header_color)
    head.line.color.rgb = RGBColor(*header_color)
    
    txH = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.5))
    pH = txH.text_frame.paragraphs[0]
    pH.text = table_name
    pH.font.bold = True
    pH.font.size = Pt(16)
    pH.font.color.rgb = RGBColor(255, 255, 255)
    pH.alignment = PP_ALIGN.CENTER
    
    # Columns
    txB = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5), Inches(width - 0.2), Inches(height - 0.5))
    for col in columns:
        pC = txB.text_frame.add_paragraph()
        if "PK" in col or "FK" in col:
            pC.text = col
            pC.font.bold = True
        else:
            pC.text = col
        pC.font.size = Pt(14)
        pC.font.color.rgb = RGBColor(44, 62, 80)
        pC.space_before = Pt(2)
        
    # Return connection points
    return {
        "top": (x + width/2, y),
        "bottom": (x + width/2, y + height),
        "left": (x, y + height/2),
        "right": (x + width, y + height/2)
    }

def connect_tables(pt1, pt2, style=MSO_CONNECTOR.ELBOW):
    conn = slide.shapes.add_connector(
        style, 
        Inches(pt1[0]), Inches(pt1[1]), 
        Inches(pt2[0]), Inches(pt2[1])
    )
    conn.line.width = Pt(1.5)
    conn.line.color.rgb = RGBColor(149, 165, 166)
    # PPTX doesn't easily support distinct 1:N crow's feet via basic API, 
    # but the elbow lines provide clear visual structure.

# Table definitions
blue = (41, 128, 185)
orange = (211, 84, 0)
green = (39, 174, 96)
red = (192, 57, 43)
purple = (142, 68, 173)

# 1. Users (Center-Left)
t_users = draw_table(
    1.0, 3.0, "👤 Users",
    ["🔑 id (PK)", "📝 name", "📏 age, height"],
    blue
)

# 2. Nutrition Plans (Top-Right)
t_nutr = draw_table(
    5.5, 1.2, "🍎 Nutrition_Plans",
    ["🔑 id (PK)", "🔗 user_id (FK)", "🔥 daily_calories"],
    orange
)

# 3. Meals (Far-Right Top)
t_meals = draw_table(
    9.5, 1.2, "🥗 Meals",
    ["🔑 id (PK)", "🔗 plan_id (FK)", "🥑 food_name"],
    orange
)

# 4. Workouts (Center-Right)
t_work = draw_table(
    5.5, 3.4, "🏋️ Workouts",
    ["🔑 id (PK)", "🔗 user_id (FK)", "📅 workout_type"],
    green
)

# 5. Exercises (Far-Right Middle)
t_ex = draw_table(
    9.5, 3.4, "🏃‍♂️ Exercises",
    ["🔑 id (PK)", "🔗 workout_id (FK)", "🎯 exercise_name"],
    green
)

# 6. CV Logs (Bottom-Right)
t_cv = draw_table(
    5.5, 5.6, "📷 CV_Logs",
    ["🔑 id (PK)", "🔗 user_id (FK)", "✅ accuracy_score"],
    purple
)

# 7. Motivation (Bottom-Left)
t_motiv = draw_table(
    1.0, 5.6, "🔥 Motivation_Streak",
    ["🔑 id (PK)", "🔗 user_id (FK)", "⭐ current_streak"],
    red
)

# Draw lines
connect_tables(t_users["right"], t_nutr["left"])
connect_tables(t_nutr["right"], t_meals["left"], MSO_CONNECTOR.STRAIGHT)
connect_tables(t_users["right"], t_work["left"])
connect_tables(t_work["right"], t_ex["left"], MSO_CONNECTOR.STRAIGHT)
connect_tables(t_users["right"], t_cv["left"])
connect_tables(t_users["bottom"], t_motiv["top"], MSO_CONNECTOR.STRAIGHT)

# ============================================================
# 3. ТӨМЕНДЕГІ ТҮСІНІКТЕМЕ
# ============================================================
txFoot = slide.shapes.add_textbox(Inches(5.0), Inches(5.8), Inches(7.5), Inches(1.2))
tfFoot = txFoot.text_frame
tfFoot.word_wrap = True
pFoot = tfFoot.paragraphs[0]
pFoot.text = (
    "Фитнес деректері қатаң байланыстарды (Relationships) талап етеді. "
    "Бір User-де бірнеше Workout, бір Workout-та бірнеше Exercise болады. "
    "Сондықтан біз NoSQL емес, Neon (PostgreSQL) реляциялық базасын таңдадық."
)
pFoot.font.size = Pt(16)
pFoot.font.italic = True
pFoot.font.color.rgb = RGBColor(127, 140, 141)

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
