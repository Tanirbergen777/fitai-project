from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "inference_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
pptx_path = OUTPUT_DIR / "lunge_model_slide.pptx"
img_path = OUTPUT_DIR / "lunge_cm_presentation.png"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 1. Title
txBox_title = slide.shapes.add_textbox(Inches(0), Inches(0.4), Inches(13.33), Inches(1))
tf_title = txBox_title.text_frame
p_title = tf_title.paragraphs[0]
p_title.text = "Lunge моделі — Зерттеу шектеулері (Limitations)"
p_title.font.bold = True
p_title.font.size = Pt(32)
p_title.font.color.rgb = RGBColor(192, 57, 43) # Dark Red for Limitation
p_title.alignment = PP_ALIGN.CENTER

# 2. Top description
txBox_desc = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.33), Inches(1))
tf_desc = txBox_desc.text_frame
tf_desc.word_wrap = True
p_desc = tf_desc.paragraphs[0]
p_desc.text = "Ғылыми жұмыста тек жетістіктерді ғана емес, шектеулерді (limitations) де көрсету маңызды. Lunge (алға қадам басу) жаттығуын бағалау моделі біз үшін қиындық тудырды."
p_desc.font.size = Pt(18)
p_desc.font.color.rgb = RGBColor(52, 73, 94)

# 3. Add Image
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(2.6), height=Inches(3.8))

# 4. Results text box
txBox_res = slide.shapes.add_textbox(Inches(6.5), Inches(2.6), Inches(6), Inches(3.8))
tf_res = txBox_res.text_frame
tf_res.word_wrap = True

def add_bullet(tf, bold_text, normal_text, highlight=False):
    p = tf.add_paragraph()
    p.text = bold_text
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(192, 57, 43) if highlight else RGBColor(41, 128, 185)
    
    run = p.add_run()
    run.text = normal_text
    run.font.bold = False
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(52, 73, 94)
    p.space_after = Pt(14)

tf_res.text = "" # clear default paragraph
add_bullet(tf_res, "• Алгоритм: ", "Random Forest")
add_bullet(tf_res, "• Дәлдік (Accuracy): ", "Бар болғаны 60.0%", highlight=True)
add_bullet(tf_res, "• Неге төмен болды?: ", "Бұл жаттығуға арналған датасет өте аз болды (тек 174 қатар жиналды). Алдыңғы модельдерде 3000-5000 қатар болған.")
add_bullet(tf_res, "• Шешу жолы (Болашақта): ", "Дәлдікті көтеру үшін деректер қорын (dataset) ұлғайту немесе 3D камералар (Lidar) қолдану қажет.")

# 5. Conclusion (Footer)
txBox_foot = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.8))
tf_foot = txBox_foot.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.text = "Қорытынды: Модельдің сапасы тікелей деректердің (Data) көлеміне тәуелді екенін осы тәжірибеден анық көруге болады."
p_foot.font.bold = True
p_foot.font.italic = True
p_foot.font.size = Pt(18)
p_foot.font.color.rgb = RGBColor(142, 68, 173) # Purple for learning
p_foot.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(f"Презентация сақталды: {pptx_path}")
